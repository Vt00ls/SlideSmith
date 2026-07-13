#!/usr/bin/env python3

from __future__ import annotations

import argparse
import contextlib
import hashlib
import io
import json
import os
import re
import shutil
import stat
import subprocess
import sys
import tempfile
import unittest
import zipfile
from pathlib import Path
from unittest import mock


sys.path.insert(0, str(Path(__file__).resolve().parent))
import ppt_runner  # noqa: E402


CONTENT_TYPES = {
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
    ".pptm": "application/vnd.ms-powerpoint.presentation.macroEnabled.main+xml",
    ".ppsx": "application/vnd.openxmlformats-officedocument.presentationml.slideshow.main+xml",
    ".ppsm": "application/vnd.ms-powerpoint.slideshow.macroEnabled.main+xml",
    ".potx": "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
    ".potm": "application/vnd.ms-powerpoint.template.macroEnabled.main+xml",
}

CONTENT_TYPES_NAMESPACE = "http://schemas.openxmlformats.org/package/2006/content-types"
PRESENTATIONML_NAMESPACE = "http://schemas.openxmlformats.org/presentationml/2006/main"

RUNNER_PATH = Path(__file__).resolve().with_name("ppt_runner.py")


class IsolatedRunnerStateTestCase(unittest.TestCase):
    def setUp(self) -> None:
        super().setUp()
        self.runner_state_tmp = tempfile.TemporaryDirectory()
        self.addCleanup(self.runner_state_tmp.cleanup)
        self.workspace = Path(self.runner_state_tmp.name) / "workspace"
        self.workspace.mkdir(parents=True)
        state_dir = self.workspace / ".slidesmith"
        self.runner_globals = mock.patch.multiple(
            ppt_runner,
            WORKSPACE=self.workspace,
            STATE_DIR=state_dir,
            EVENTS_PATH=state_dir / "events.ndjson",
            STATUS_PATH=state_dir / "status.json",
            ARTIFACTS_PATH=state_dir / "artifacts.json",
        )
        self.runner_globals.start()
        self.addCleanup(self.runner_globals.stop)


class PPTSourceStagingTests(IsolatedRunnerStateTestCase):
    def test_stage_normalizes_slideshow_and_template_main_content_types(self) -> None:
        expected = {
            ".ppsx": CONTENT_TYPES[".pptx"],
            ".ppsm": CONTENT_TYPES[".pptm"],
            ".potx": CONTENT_TYPES[".pptx"],
            ".potm": CONTENT_TYPES[".pptm"],
        }
        for extension, expected_content_type in expected.items():
            with self.subTest(extension=extension), tempfile.TemporaryDirectory() as tmp:
                root = Path(tmp)
                source = root / f"deck{extension}"
                write_ooxml_package(source, CONTENT_TYPES[extension])

                staged = stage_one(source, root / "staged")

                self.assertEqual(read_main_content_type(staged), expected_content_type)
                self.assertEqual(read_zip_entry(staged, "ppt/marker.bin"), b"preserve-me")
                self.assertEqual(staged.name, source.name)

    def test_stage_keeps_pptx_and_pptm_byte_for_byte(self) -> None:
        for extension in (".pptx", ".pptm"):
            with self.subTest(extension=extension), tempfile.TemporaryDirectory() as tmp:
                root = Path(tmp)
                source = root / f"deck{extension}"
                write_ooxml_package(source, CONTENT_TYPES[extension])
                original = source.read_bytes()

                staged = stage_one(source, root / "staged")

                self.assertEqual(staged.read_bytes(), original)

    def test_stage_rejects_malformed_byte_preserved_and_normalized_packages(self) -> None:
        for extension in (".pptx", ".ppsx"):
            with self.subTest(extension=extension), tempfile.TemporaryDirectory() as tmp:
                root = Path(tmp)
                source = root / f"broken{extension}"
                source.write_bytes(b"not a zip package")

                with self.assertRaisesRegex(ValueError, "malformed OOXML presentation package"):
                    stage_one(source, root / "staged")

    def test_stage_rejects_missing_expected_main_content_type(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "wrong.ppsx"
            write_ooxml_package(source, CONTENT_TYPES[".pptx"])

            with self.assertRaisesRegex(ValueError, "expected slideshow main content type"):
                stage_one(source, root / "staged")

    def test_stage_rejects_missing_presentation_part_for_byte_preserved_and_normalized_packages(self) -> None:
        for extension in (".pptx", ".potx"):
            with self.subTest(extension=extension), tempfile.TemporaryDirectory() as tmp:
                root = Path(tmp)
                source = root / f"missing-part{extension}"
                write_ooxml_package(source, CONTENT_TYPES[extension], include_presentation=False)

                with self.assertRaisesRegex(ValueError, "missing ppt/presentation.xml"):
                    stage_one(source, root / "staged")

    def test_stage_rejects_non_presentationml_root_for_byte_preserved_and_normalized_packages(self) -> None:
        invalid_root = b"<p:presentation xmlns:p='urn:not-presentationml'/>"
        for extension in (".pptm", ".potm"):
            with self.subTest(extension=extension), tempfile.TemporaryDirectory() as tmp:
                root = Path(tmp)
                source = root / f"wrong-root{extension}"
                write_ooxml_package(source, CONTENT_TYPES[extension], presentation_xml=invalid_root)

                with self.assertRaisesRegex(ValueError, "PresentationML presentation root"):
                    stage_one(source, root / "staged")

    def test_stage_rejects_invalid_content_types_namespace_and_duplicate_main_override(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "wrong-namespace.pptx"
            write_ooxml_package(
                source,
                CONTENT_TYPES[".pptx"],
                content_types_namespace="urn:not-opc-content-types",
            )
            with self.assertRaisesRegex(ValueError, "content-types namespace"):
                stage_one(source, root / "staged-namespace")

            duplicate = root / "duplicate.ppsm"
            write_ooxml_package(
                duplicate,
                CONTENT_TYPES[".ppsm"],
                duplicate_main_override=True,
            )
            with self.assertRaisesRegex(ValueError, "exactly one.*main content type"):
                stage_one(duplicate, root / "staged-duplicate")

    def test_text_sources_are_readable_by_real_lite_reader(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            project = Path(tmp) / "project"
            source = project / "sources" / "notes.text"
            source.parent.mkdir(parents=True)
            source.write_text("# Notes\nReadable downstream content.\n", encoding="utf-8")

            docs = ppt_runner.load_source_documents(project)

            self.assertEqual(len(docs), 1)
            self.assertEqual(docs[0]["path"], "sources/notes.text")
            self.assertIn("Readable downstream content", docs[0]["text"])


class TemplateFillRunnerTests(IsolatedRunnerStateTestCase):
    def test_discover_template_fill_inputs_returns_every_runtime_path(self) -> None:
        project = make_template_fill_project(Path(self.runner_state_tmp.name) / "brand_project")
        (project / "sources" / "brand.md").write_text(
            "# Generated template readback\n",
            encoding="utf-8",
        )

        inputs = ppt_runner.discover_template_fill_inputs(project)

        self.assertEqual(inputs["project_path"], project)
        self.assertEqual(inputs["source_pptx"], project / "sources" / "brand.pptx")
        self.assertEqual(inputs["slide_library"], project / "analysis" / "brand.slide_library.json")
        self.assertEqual(inputs["fill_plan"], project / "analysis" / "fill_plan.json")
        self.assertEqual(inputs["check_report"], project / "analysis" / "check_report.json")
        self.assertEqual(inputs["validate_report"], project / "validation" / "validate_report.json")
        self.assertEqual(inputs["readback"], project / "validation" / "readback.md")
        self.assertEqual(
            inputs["export_base"],
            project / "exports" / "brand_project_template_fill.pptx",
        )
        self.assertEqual(inputs["content_sources"], [project / "sources" / "content.md"])

    def test_discover_template_fill_inputs_sorts_supported_content_case_insensitively(self) -> None:
        project = make_template_fill_project(Path(self.runner_state_tmp.name) / "content_project")
        (project / "sources" / "content.md").unlink()
        for name in ("zeta.TSV", "alpha.markdown", "middle.TEXT"):
            (project / "sources" / name).write_text("content\n", encoding="utf-8")

        inputs = ppt_runner.discover_template_fill_inputs(project)

        self.assertEqual(
            inputs["content_sources"],
            [
                project / "sources" / "alpha.markdown",
                project / "sources" / "middle.TEXT",
                project / "sources" / "zeta.TSV",
            ],
        )

    def test_discover_template_fill_inputs_preserves_explicit_same_stem_markdown(self) -> None:
        project = make_template_fill_project(self.workspace / "projects" / "brand_project")
        (project / "sources" / "content.md").unlink()
        (project / "sources" / "brand.md").write_text("# Business content\n", encoding="utf-8")
        write_json_file(
            self.workspace / ".slidesmith" / "source_inputs.json",
            {
                "schema": "slidesmith.source_inputs.v1",
                "files": [
                    {"name": "brand.pptx", "upload_path": "uploads/task/brand.pptx"},
                    {"name": "brand.md", "upload_path": "uploads/task/brand.md"},
                ],
            },
        )

        inputs = ppt_runner.discover_template_fill_inputs(project)

        self.assertEqual(inputs["content_sources"], [project / "sources" / "brand.md"])

    def test_discover_template_fill_inputs_rejects_malformed_matching_manifest_entry(self) -> None:
        project = make_template_fill_project(self.workspace / "projects" / "brand_project")
        (project / "sources" / "content.md").unlink()
        (project / "sources" / "brand.md").write_text("# Generated readback\n", encoding="utf-8")
        write_json_file(
            self.workspace / ".slidesmith" / "source_inputs.json",
            {
                "schema": "slidesmith.source_inputs.v1",
                "files": [{"name": "brand.md", "upload_path": 7}],
            },
        )

        with self.assertRaisesRegex(ValueError, "upload_path must be a string"):
            ppt_runner.discover_template_fill_inputs(project)

    def test_discover_template_fill_inputs_validates_entries_after_provenance_match(self) -> None:
        project = make_template_fill_project(self.workspace / "projects" / "brand_project")
        (project / "sources" / "content.md").unlink()
        (project / "sources" / "brand.md").write_text("# Explicit business content\n", encoding="utf-8")
        write_json_file(
            self.workspace / ".slidesmith" / "source_inputs.json",
            {
                "schema": "slidesmith.source_inputs.v1",
                "files": [
                    {"name": "brand.md", "upload_path": "uploads/task/brand.md"},
                    {"name": "later.md", "upload_path": 7},
                ],
            },
        )

        with self.assertRaisesRegex(ValueError, r"files\[1\]\.upload_path must be a string"):
            ppt_runner.discover_template_fill_inputs(project)

    def test_discover_template_fill_inputs_rejects_other_or_multiple_presentations_deterministically(self) -> None:
        cases = {
            "other format": ["template.PPTM"],
            "pptx plus other": ["a.pptx", "b.potm"],
            "multiple pptx": ["z.pptx", "a.PPTX"],
        }
        for name, presentations in cases.items():
            with self.subTest(name=name), tempfile.TemporaryDirectory() as tmp:
                project = Path(tmp) / "project"
                (project / "sources").mkdir(parents=True)
                (project / "sources" / "content.md").write_text("# Content\n", encoding="utf-8")
                for filename in presentations:
                    (project / "sources" / filename).write_text("presentation", encoding="utf-8")

                with self.assertRaisesRegex(ValueError, "requires exactly one source PPTX") as caught:
                    ppt_runner.discover_template_fill_inputs(project)

                message = str(caught.exception)
                self.assertIn(f"found {len(presentations)} presentation files", message)
                listed = [f"sources/{filename}" for filename in sorted(presentations)]
                self.assertEqual([message.index(path) for path in listed], sorted(message.index(path) for path in listed))

    def test_discover_template_fill_inputs_requires_business_content(self) -> None:
        project = make_template_fill_project(Path(self.runner_state_tmp.name) / "project")
        (project / "sources" / "content.md").unlink()
        (project / "sources" / "brand.md").write_text("# Generated readback\n", encoding="utf-8")
        (project / "sources" / "archived.xls").write_text("archived", encoding="utf-8")

        with self.assertRaisesRegex(ValueError, "requires content source"):
            ppt_runner.discover_template_fill_inputs(project)

    def test_discover_template_fill_inputs_requires_nonempty_matching_library(self) -> None:
        for name, contents in (("missing", None), ("empty", "")):
            with self.subTest(name=name), tempfile.TemporaryDirectory() as tmp:
                project = make_template_fill_project(Path(tmp) / "project")
                library = project / "analysis" / "brand.slide_library.json"
                library.unlink()
                if contents is not None:
                    library.write_text(contents, encoding="utf-8")

                with self.assertRaisesRegex(ValueError, "requires slide library"):
                    ppt_runner.discover_template_fill_inputs(project)

    def test_discover_template_fill_inputs_rejects_symlinked_content(self) -> None:
        project = make_template_fill_project(Path(self.runner_state_tmp.name) / "project")
        (project / "sources" / "content.md").unlink()
        (project / "sources" / "actual.md").write_text("# Content\n", encoding="utf-8")
        (project / "sources" / "content.md").symlink_to("actual.md")

        with self.assertRaisesRegex(ValueError, "non-symlinked"):
            ppt_runner.discover_template_fill_inputs(project)

    def test_template_fill_check_accepts_only_exit_zero_or_one_with_fresh_valid_report(self) -> None:
        for returncode, summary in (
            (0, {"ok": 3, "warn": 1, "error": 0}),
            (1, {"ok": 2, "warn": 0, "error": 1}),
        ):
            with self.subTest(returncode=returncode), tempfile.TemporaryDirectory() as tmp:
                project = make_template_fill_project(Path(tmp) / "project")
                report_path = project / "analysis" / "check_report.json"

                def fake_run(command: list[str], **kwargs: object) -> subprocess.CompletedProcess[str]:
                    write_template_fill_report(report_path, "template_fill_pptx_check.v1", summary)
                    return subprocess.CompletedProcess(command, returncode, "", "")

                output = io.StringIO()
                with mock.patch.object(ppt_runner, "run_command", side_effect=fake_run) as run, contextlib.redirect_stdout(output):
                    result = ppt_runner.template_fill_check(arg_namespace(project))

                self.assertEqual(result, summary)
                command = run.call_args.args[0]
                self.assertEqual(command[2], "check-plan")
                self.assertEqual(command[-2:], ["-o", str(report_path)])
                self.assertIs(run.call_args.kwargs["check"], False)
                self.assertIn(f"ok={summary['ok']} warn={summary['warn']} error={summary['error']}", output.getvalue())
                events = read_events()
                check_event = [event for event in events if event["type"] == "template_fill_check"][-1]
                self.assertEqual(check_event["payload"]["summary"], summary)

    def test_template_fill_check_removes_stale_report_before_invocation(self) -> None:
        project = make_template_fill_project(Path(self.runner_state_tmp.name) / "project")
        report_path = project / "analysis" / "check_report.json"
        write_template_fill_report(
            report_path,
            "template_fill_pptx_check.v1",
            {"ok": 99, "warn": 0, "error": 0},
        )

        def fake_run(command: list[str], **kwargs: object) -> subprocess.CompletedProcess[str]:
            self.assertFalse(report_path.exists(), "stale check report survived until invocation")
            return subprocess.CompletedProcess(command, 1, "", "")

        with mock.patch.object(ppt_runner, "run_command", side_effect=fake_run):
            with self.assertRaisesRegex(RuntimeError, "check report"):
                ppt_runner.template_fill_check(arg_namespace(project))
        self.assertFalse(report_path.exists())

    def test_template_fill_check_rejects_invalid_or_system_failure_reports(self) -> None:
        cases = (
            ("missing", 0, None),
            ("corrupt", 0, "corrupt"),
            ("wrong schema", 1, "wrong_schema"),
            ("invalid summary", 1, "invalid_summary"),
            ("system exit", 2, "valid"),
        )
        for name, returncode, report_kind in cases:
            with self.subTest(name=name), tempfile.TemporaryDirectory() as tmp:
                project = make_template_fill_project(Path(tmp) / "project")
                report_path = project / "analysis" / "check_report.json"

                def fake_run(command: list[str], **kwargs: object) -> subprocess.CompletedProcess[str]:
                    if report_kind == "corrupt":
                        report_path.write_text("{", encoding="utf-8")
                    elif report_kind == "wrong_schema":
                        write_template_fill_report(report_path, "wrong.v1", {"ok": 1, "warn": 0, "error": 0})
                    elif report_kind == "invalid_summary":
                        write_template_fill_report(report_path, "template_fill_pptx_check.v1", {"ok": 1, "warn": "0", "error": 0})
                    elif report_kind == "valid":
                        write_template_fill_report(report_path, "template_fill_pptx_check.v1", {"ok": 1, "warn": 0, "error": 0})
                    return subprocess.CompletedProcess(command, returncode, "", "system failure")

                with mock.patch.object(ppt_runner, "run_command", side_effect=fake_run):
                    with self.assertRaises((RuntimeError, ValueError, json.JSONDecodeError)):
                        ppt_runner.template_fill_check(arg_namespace(project))

    def test_template_fill_apply_passes_transition_without_force_and_requires_new_timestamped_export(self) -> None:
        project = make_template_fill_project(Path(self.runner_state_tmp.name) / "apply_project")
        stale = project / "exports" / "apply_project_template_fill_20260712_120000.pptx"
        stale.parent.mkdir(parents=True)
        stale.write_text("old", encoding="utf-8")
        fresh = project / "exports" / "apply_project_template_fill_20260713_120000.pptx"

        def fake_run(command: list[str], **kwargs: object) -> subprocess.CompletedProcess[str]:
            fresh.write_text("new", encoding="utf-8")
            return subprocess.CompletedProcess(command, 0, "", "")

        with mock.patch.object(ppt_runner, "run_command", side_effect=fake_run) as run:
            result = ppt_runner.template_fill_apply(arg_namespace(project, transition="push"))

        self.assertEqual(result, fresh)
        command = run.call_args.args[0]
        self.assertEqual(command[2], "apply")
        self.assertEqual(command[-2:], ["--transition", "push"])
        self.assertNotIn("--force", command)
        self.assertIs(run.call_args.kwargs["check"], False)

    def test_template_fill_apply_rejects_nonzero_and_main_writes_failed_status(self) -> None:
        project = make_template_fill_project(Path(self.runner_state_tmp.name) / "apply_project")

        def fake_run(command: list[str], **kwargs: object) -> subprocess.CompletedProcess[str]:
            return subprocess.CompletedProcess(command, 1, "", "apply failed")

        argv = ["ppt_runner.py", "template-fill-apply", "--project-path", str(project)]
        with mock.patch.object(ppt_runner, "run_command", side_effect=fake_run), mock.patch.object(sys, "argv", argv):
            with self.assertRaisesRegex(RuntimeError, "exit 1"):
                ppt_runner.main()

        status = json.loads(ppt_runner.STATUS_PATH.read_text(encoding="utf-8"))
        self.assertEqual(status["status"], "failed")
        self.assertIn("apply", status["error"])

    def test_template_fill_apply_rejects_stale_or_untimestamped_exports(self) -> None:
        for name, create_untimestamped in (("stale", False), ("untimestamped", True)):
            with self.subTest(name=name), tempfile.TemporaryDirectory() as tmp:
                project = make_template_fill_project(Path(tmp) / "apply_project")
                exports = project / "exports"
                exports.mkdir(parents=True)
                (exports / "apply_project_template_fill_20260712_120000.pptx").write_text("old", encoding="utf-8")

                def fake_run(command: list[str], **kwargs: object) -> subprocess.CompletedProcess[str]:
                    if create_untimestamped:
                        (exports / "apply_project_template_fill.pptx").write_text("not timestamped", encoding="utf-8")
                    return subprocess.CompletedProcess(command, 0, "", "")

                with mock.patch.object(ppt_runner, "run_command", side_effect=fake_run):
                    with self.assertRaisesRegex(RuntimeError, "new timestamped PPTX"):
                        ppt_runner.template_fill_apply(arg_namespace(project, transition="fade"))

    def test_template_fill_validate_clears_stale_outputs_and_allows_warnings(self) -> None:
        project = make_template_fill_project(Path(self.runner_state_tmp.name) / "validate_project")
        readback = project / "validation" / "readback.md"
        report_path = project / "validation" / "validate_report.json"
        readback.parent.mkdir(parents=True)
        readback.write_text("stale readback", encoding="utf-8")
        write_template_fill_report(report_path, "template_fill_pptx_validate.v1", {"ok": 99, "warn": 0, "error": 0})
        summary = {"ok": 4, "warn": 2, "error": 0}

        def fake_run(command: list[str], **kwargs: object) -> subprocess.CompletedProcess[str]:
            self.assertFalse(readback.exists(), "stale readback survived until invocation")
            self.assertFalse(report_path.exists(), "stale validation report survived until invocation")
            readback.write_text("## Slide 1\n", encoding="utf-8")
            write_template_fill_report(report_path, "template_fill_pptx_validate.v1", summary)
            return subprocess.CompletedProcess(command, 0, "", "")

        output = io.StringIO()
        with mock.patch.object(ppt_runner, "run_command", side_effect=fake_run) as run, contextlib.redirect_stdout(output):
            result = ppt_runner.template_fill_validate(arg_namespace(project))

        self.assertEqual(result, summary)
        self.assertEqual(run.call_args.args[0][2:], ["validate", str(project)])
        self.assertIs(run.call_args.kwargs["check"], False)
        self.assertIn("ok=4 warn=2 error=0", output.getvalue())

    def test_template_fill_validate_rejects_exit_errors_and_invalid_outputs(self) -> None:
        cases = (
            ("nonzero", 1, "valid", "valid"),
            ("report errors", 0, "errors", "valid"),
            ("missing report", 0, None, "valid"),
            ("wrong schema", 0, "wrong_schema", "valid"),
            ("corrupt report", 0, "corrupt", "valid"),
            ("missing readback", 0, "valid", None),
            ("empty readback", 0, "valid", "empty"),
        )
        for name, returncode, report_kind, readback_kind in cases:
            with self.subTest(name=name), tempfile.TemporaryDirectory() as tmp:
                project = make_template_fill_project(Path(tmp) / "validate_project")
                readback = project / "validation" / "readback.md"
                report_path = project / "validation" / "validate_report.json"

                def fake_run(command: list[str], **kwargs: object) -> subprocess.CompletedProcess[str]:
                    if readback_kind is not None:
                        readback.parent.mkdir(parents=True, exist_ok=True)
                        readback.write_text("" if readback_kind == "empty" else "## Slide 1\n", encoding="utf-8")
                    if report_kind == "valid":
                        write_template_fill_report(report_path, "template_fill_pptx_validate.v1", {"ok": 1, "warn": 0, "error": 0})
                    elif report_kind == "errors":
                        write_template_fill_report(report_path, "template_fill_pptx_validate.v1", {"ok": 0, "warn": 0, "error": 1})
                    elif report_kind == "wrong_schema":
                        write_template_fill_report(report_path, "wrong.v1", {"ok": 1, "warn": 0, "error": 0})
                    elif report_kind == "corrupt":
                        report_path.parent.mkdir(parents=True, exist_ok=True)
                        report_path.write_text("{", encoding="utf-8")
                    return subprocess.CompletedProcess(command, returncode, "", "validation failed")

                with mock.patch.object(ppt_runner, "run_command", side_effect=fake_run):
                    with self.assertRaises((RuntimeError, ValueError, json.JSONDecodeError)):
                        ppt_runner.template_fill_validate(arg_namespace(project))

    def test_publish_project_includes_validation_and_uses_isolated_state(self) -> None:
        project = Path(self.runner_state_tmp.name) / "publish_project"
        (project / "exports").mkdir(parents=True)
        (project / "validation").mkdir(parents=True)
        (project / "exports" / "result.pptx").write_text("pptx", encoding="utf-8")
        (project / "validation" / "readback.md").write_text("## Slide 1\n", encoding="utf-8")

        ppt_runner.publish_project(project)

        manifest = json.loads((project / ".slidesmith-artifacts.json").read_text(encoding="utf-8"))
        paths = {item["path"] for item in manifest["artifacts"]}
        self.assertIn("exports/result.pptx", paths)
        self.assertIn("validation/readback.md", paths)
        self.assertTrue(ppt_runner.ARTIFACTS_PATH.is_relative_to(self.workspace))
        self.assertTrue(ppt_runner.ARTIFACTS_PATH.is_file())


def make_template_fill_project(project: Path) -> Path:
    project = project.resolve()
    (project / "sources").mkdir(parents=True)
    (project / "analysis").mkdir(parents=True)
    (project / "sources" / "brand.pptx").write_text("pptx", encoding="utf-8")
    (project / "sources" / "content.md").write_text("# Content\n", encoding="utf-8")
    (project / "analysis" / "brand.slide_library.json").write_text(
        '{"slides":[]}\n',
        encoding="utf-8",
    )
    return project


def arg_namespace(project: Path, *, transition: str = "fade") -> argparse.Namespace:
    return argparse.Namespace(project="", project_path=str(project), transition=transition)


def write_json_file(path: Path, payload: object) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, indent=2) + "\n", encoding="utf-8")


def write_template_fill_report(path: Path, schema: str, summary: dict[str, object]) -> None:
    write_json_file(path, {"schema": schema, "summary": summary, "results": []})


def read_events() -> list[dict[str, object]]:
    return [
        json.loads(line)
        for line in ppt_runner.EVENTS_PATH.read_text(encoding="utf-8").splitlines()
        if line.strip()
    ]


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def snapshot_filesystem_node(path: Path) -> tuple[str, ...]:
    try:
        metadata = path.lstat()
    except FileNotFoundError:
        return ("missing",)
    if stat.S_ISLNK(metadata.st_mode):
        return ("symlink", os.readlink(path))
    if stat.S_ISREG(metadata.st_mode):
        return ("file", str(metadata.st_size), sha256_file(path))
    if stat.S_ISDIR(metadata.st_mode):
        return ("directory",)
    return ("other", oct(stat.S_IFMT(metadata.st_mode)), str(metadata.st_size))


def snapshot_tree_manifest(root: Path) -> tuple[tuple[str, tuple[str, ...]], ...]:
    root_snapshot = snapshot_filesystem_node(root)
    entries: list[tuple[str, tuple[str, ...]]] = [(".", root_snapshot)]
    if root_snapshot[0] != "directory":
        return tuple(entries)

    for current, directory_names, file_names in os.walk(root, topdown=True, followlinks=False):
        directory_names.sort()
        file_names.sort()
        current_path = Path(current)
        for name in [*directory_names, *file_names]:
            path = current_path / name
            entries.append((path.relative_to(root).as_posix(), snapshot_filesystem_node(path)))
    return tuple(entries)


def snapshot_external_pycache_manifest(
    external_root: Path,
) -> tuple[tuple[str, tuple[str, ...]], ...]:
    entries: list[tuple[str, tuple[str, ...]]] = []
    skill_root = external_root / "skills" / "ppt-master"
    if snapshot_filesystem_node(skill_root)[0] != "directory":
        return tuple(entries)

    for current, directory_names, _ in os.walk(skill_root, topdown=True, followlinks=False):
        directory_names[:] = sorted(name for name in directory_names if name != ".git")
        current_path = Path(current)
        cache_names = [name for name in directory_names if name == "__pycache__"]
        for name in cache_names:
            cache = current_path / name
            cache_relative = cache.relative_to(external_root)
            for relative, node in snapshot_tree_manifest(cache):
                path = cache_relative if relative == "." else cache_relative / relative
                entries.append((path.as_posix(), node))
        directory_names[:] = [name for name in directory_names if name != "__pycache__"]
    return tuple(sorted(entries, key=lambda item: item[0]))


def snapshot_external_skill_tree_manifest(
    external_root: Path,
) -> tuple[tuple[str, tuple[str, ...]], ...]:
    skill_root = external_root / "skills" / "ppt-master"
    root_snapshot = snapshot_filesystem_node(skill_root)
    entries: list[tuple[str, tuple[str, ...]]] = [(".", root_snapshot)]
    if root_snapshot[0] != "directory":
        return tuple(entries)

    for current, directory_names, file_names in os.walk(skill_root, topdown=True, followlinks=False):
        directory_names[:] = sorted(name for name in directory_names if name != ".git")
        file_names = sorted(name for name in file_names if name != ".git")
        current_path = Path(current)
        for name in [*directory_names, *file_names]:
            path = current_path / name
            entries.append((path.relative_to(skill_root).as_posix(), snapshot_filesystem_node(path)))
    return tuple(entries)


def snapshot_git_status(root: Path) -> tuple[int, str, str]:
    try:
        completed = subprocess.run(
            ["git", "-C", str(root), "status", "--porcelain=v1", "--untracked-files=all"],
            env={**os.environ, "GIT_OPTIONAL_LOCKS": "0"},
            text=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
        )
    except OSError as exc:
        return (-1, "", f"{type(exc).__name__}: {exc}")
    return (completed.returncode, completed.stdout, completed.stderr)


def capture_template_fill_safety_snapshot(
    source_paths: list[Path] | tuple[Path, ...],
    *,
    external_root: Path,
    repository_state: Path,
) -> dict[str, object]:
    sources = tuple(Path(path) for path in source_paths)
    return {
        "source_paths": sources,
        "external_root": Path(external_root),
        "repository_state": Path(repository_state),
        "source fixtures": tuple(
            (str(path), snapshot_filesystem_node(path))
            for path in sources
        ),
        "external ppt-master Git status": snapshot_git_status(external_root),
        "external ppt-master __pycache__": snapshot_external_pycache_manifest(external_root),
        "external ppt-master skill tree": snapshot_external_skill_tree_manifest(external_root),
        "repository runtime .slidesmith": snapshot_tree_manifest(repository_state),
    }


def assert_template_fill_safety_unchanged(snapshot: dict[str, object]) -> None:
    differences: list[str] = []
    component_captures = (
        (
            "source fixtures",
            lambda: tuple(
                (str(path), snapshot_filesystem_node(path))
                for path in snapshot["source_paths"]
            ),
        ),
        (
            "external ppt-master Git status",
            lambda: snapshot_git_status(snapshot["external_root"]),
        ),
        (
            "external ppt-master __pycache__",
            lambda: snapshot_external_pycache_manifest(snapshot["external_root"]),
        ),
        (
            "external ppt-master skill tree",
            lambda: snapshot_external_skill_tree_manifest(snapshot["external_root"]),
        ),
        (
            "repository runtime .slidesmith",
            lambda: snapshot_tree_manifest(snapshot["repository_state"]),
        ),
    )
    for label, capture_current in component_captures:
        try:
            current = capture_current()
        except Exception as exc:
            differences.append(f"{label} capture failed: {type(exc).__name__}: {exc}")
            continue
        if current != snapshot[label]:
            differences.append(
                f"{label} changed\n"
                f"  before: {snapshot[label]!r}\n"
                f"  after:  {current!r}"
            )
    if differences:
        raise AssertionError("Template Fill smoke safety violation:\n" + "\n".join(differences))


def register_template_fill_safety_finalizer(
    test_case: unittest.TestCase,
    source_paths: list[Path] | tuple[Path, ...],
    *,
    external_root: Path,
    repository_state: Path,
) -> dict[str, object]:
    snapshot = capture_template_fill_safety_snapshot(
        source_paths,
        external_root=external_root,
        repository_state=repository_state,
    )
    test_case.addCleanup(assert_template_fill_safety_unchanged, snapshot)
    return snapshot


def run_template_fill_smoke_subprocess(
    command: list[str],
    *,
    workspace: Path,
    temporary_root: Path,
    timeout: int = 180,
) -> subprocess.CompletedProcess[str]:
    resolved_root = temporary_root.resolve()
    resolved_workspace = workspace.resolve()
    try:
        resolved_workspace.relative_to(resolved_root)
    except ValueError as exc:
        raise ValueError(f"smoke workspace must be below temporary root: {workspace}") from exc

    controlled_cwd = temporary_root / "subprocess-cwd"
    controlled_cwd.mkdir(parents=True, exist_ok=True)
    environment = os.environ.copy()
    environment["PYTHONDONTWRITEBYTECODE"] = "1"
    environment["WORKSPACE"] = str(workspace)
    return subprocess.run(
        command,
        cwd=controlled_cwd,
        env=environment,
        text=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        timeout=timeout,
    )


def assert_template_fill_runner_state_isolated(
    temporary_root: Path,
    workspace: Path,
) -> None:
    resolved_workspace = workspace.resolve()
    state_directories = sorted(temporary_root.rglob(".slidesmith"), key=lambda path: str(path))
    escaped_paths: list[Path] = []
    for state_path in state_directories:
        try:
            state_path.resolve().relative_to(resolved_workspace)
        except ValueError:
            escaped_paths.append(state_path)
    expected_state = workspace / ".slidesmith"
    if escaped_paths:
        raise AssertionError(f"runner state outside workspace: {escaped_paths}")
    if state_directories != [expected_state]:
        raise AssertionError(
            f"runner state directories = {state_directories!r}, expected {[expected_state]!r}"
        )
    for state_file in expected_state.rglob("*"):
        try:
            state_file.resolve().relative_to(resolved_workspace)
        except ValueError as exc:
            raise AssertionError(f"runner state outside workspace: {state_file}") from exc


class RealTemplateFillSafetyGuardTests(unittest.TestCase):
    def make_safety_context(self, root: Path) -> tuple[list[Path], Path, Path]:
        external_root = root / "ppt-master"
        external_root.mkdir()
        initialized = subprocess.run(
            ["git", "init", "-q", str(external_root)],
            text=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
        )
        self.assertEqual(initialized.returncode, 0, initialized.stdout + initialized.stderr)
        (external_root / ".gitignore").write_text(
            "__pycache__/\n*.log\n.cache/\n**/projects/**\n",
            encoding="utf-8",
        )
        (external_root / "skills" / "ppt-master").mkdir(parents=True)
        fixture = external_root / "fixture.pptx"
        content = external_root / "content.md"
        fixture.write_bytes(b"fixture")
        content.write_text("# Content\n", encoding="utf-8")
        repository_state = root / "runtime" / ".slidesmith"
        return [fixture, content], external_root, repository_state

    def test_smoke_subprocess_uses_controlled_temp_cwd_and_workspace(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            temporary_root = Path(tmp)
            workspace = temporary_root / "workspace"
            workspace.mkdir()
            probe = (
                "import json, os; from pathlib import Path; "
                "state = Path('.slidesmith'); state.mkdir(); "
                "(state / 'cwd-probe').write_text('probe', encoding='utf-8'); "
                "print(json.dumps({'cwd': str(Path.cwd()), "
                "'workspace': os.environ.get('WORKSPACE'), "
                "'dont_write_bytecode': os.environ.get('PYTHONDONTWRITEBYTECODE')}))"
            )
            result = run_template_fill_smoke_subprocess(
                [sys.executable, "-c", probe],
                workspace=workspace,
                temporary_root=temporary_root,
            )

            self.assertEqual(result.returncode, 0, result.stdout + result.stderr)
            payload = json.loads(result.stdout)
            controlled_cwd = temporary_root / "subprocess-cwd"
            self.assertTrue(controlled_cwd.is_dir())
            self.assertEqual(Path(payload["cwd"]), controlled_cwd.resolve())
            self.assertEqual(payload["workspace"], str(workspace))
            self.assertEqual(payload["dont_write_bytecode"], "1")
            cwd_state = controlled_cwd / ".slidesmith" / "cwd-probe"
            self.assertTrue(cwd_state.is_file())
            self.assertFalse(cwd_state.resolve().is_relative_to(workspace.resolve()))
            workspace_state = workspace / ".slidesmith"
            workspace_state.mkdir()
            (workspace_state / "status.json").write_text("{}\n", encoding="utf-8")
            with self.assertRaisesRegex(AssertionError, "runner state outside workspace"):
                assert_template_fill_runner_state_isolated(temporary_root, workspace)

    def test_safety_guard_detects_repo_state_git_status_and_pycache_changes(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            sources, external_root, repository_state = self.make_safety_context(root)
            repository_state.mkdir(parents=True)
            leaked_state = repository_state / "status.json"
            leaked_state.write_bytes(b'{"a":1}\n')
            state_stat = leaked_state.stat()
            pycache = external_root / "skills" / "ppt-master" / "scripts" / "__pycache__"
            pycache.mkdir(parents=True)
            cached_module = pycache / "module.pyc"
            cached_module.write_bytes(b"old bytecode")
            cache_stat = cached_module.stat()
            snapshot = capture_template_fill_safety_snapshot(
                sources,
                external_root=external_root,
                repository_state=repository_state,
            )

            leaked_state.write_bytes(b'{"b":1}\n')
            os.utime(leaked_state, ns=(state_stat.st_atime_ns, state_stat.st_mtime_ns))
            cached_module.write_bytes(b"new bytecode")
            os.utime(cached_module, ns=(cache_stat.st_atime_ns, cache_stat.st_mtime_ns))
            (external_root / "unexpected.txt").write_text("dirty\n", encoding="utf-8")

            with self.assertRaises(AssertionError) as caught:
                assert_template_fill_safety_unchanged(snapshot)

            message = str(caught.exception)
            self.assertIn("repository runtime .slidesmith", message)
            self.assertIn("external ppt-master Git status", message)
            self.assertIn("external ppt-master __pycache__", message)

    def test_safety_guard_detects_ignored_external_skill_tree_changes(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            sources, external_root, repository_state = self.make_safety_context(root)
            skill_root = external_root / "skills" / "ppt-master"
            (skill_root / "logs").mkdir()
            ignored_log = skill_root / "logs" / "smoke.log"
            ignored_log.write_bytes(b"before-log!")
            (skill_root / ".cache").mkdir()
            ignored_cache = skill_root / ".cache" / "state.bin"
            ignored_cache.write_bytes(b"before-cache")
            ignored_project = skill_root / "projects" / "leak" / "output.json"
            ignored_project.parent.mkdir(parents=True)
            ignored_project.write_bytes(b'{"before":1}')
            snapshot = capture_template_fill_safety_snapshot(
                sources,
                external_root=external_root,
                repository_state=repository_state,
            )
            git_status_before = snapshot_git_status(external_root)
            skill_manifest = dict(snapshot["external ppt-master skill tree"])
            self.assertEqual(skill_manifest["logs"], ("directory",))
            self.assertEqual(skill_manifest["logs/smoke.log"][:2], ("file", "11"))
            self.assertEqual(len(skill_manifest["logs/smoke.log"][2]), 64)

            ignored_log.write_bytes(b"after--log!")
            ignored_cache.write_bytes(b"after--cache")
            ignored_project.write_bytes(b'{"after-":1}')

            self.assertEqual(
                snapshot_git_status(external_root),
                git_status_before,
                "ignored writes should prove why Git status alone is insufficient",
            )
            with self.assertRaisesRegex(AssertionError, "external ppt-master skill tree"):
                assert_template_fill_safety_unchanged(snapshot)

    def test_safety_finalizer_preserves_primary_failure(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            sources, external_root, repository_state = self.make_safety_context(root)
            fixture = sources[0]

            class ForcedPhaseFailure(unittest.TestCase):
                def runTest(inner_self) -> None:
                    register_template_fill_safety_finalizer(
                        inner_self,
                        sources,
                        external_root=external_root,
                        repository_state=repository_state,
                    )
                    fixture.write_bytes(b"changed fixture")
                    inner_self.fail("forced phase failure")

            result = unittest.TestResult()
            ForcedPhaseFailure().run(result)
            diagnostics = "\n".join(
                detail
                for _, detail in [*result.failures, *result.errors]
            )
            self.assertIn("forced phase failure", diagnostics)
            self.assertIn("source fixtures", diagnostics)
            self.assertGreaterEqual(len(result.failures) + len(result.errors), 2)

    def test_safety_finalizer_exhausts_checks_after_source_capture_error(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            sources, external_root, repository_state = self.make_safety_context(root)
            fixture = sources[0]
            repository_state.mkdir(parents=True)
            state_file = repository_state / "status.json"
            state_file.write_text('{"state":"before"}\n', encoding="utf-8")
            pycache = external_root / "skills" / "ppt-master" / "scripts" / "__pycache__"
            pycache.mkdir(parents=True)
            cached_module = pycache / "module.pyc"
            cached_module.write_bytes(b"before bytecode")
            ignored_skill_log = external_root / "skills" / "ppt-master" / "runner.log"
            capture_fault = {"enabled": False}
            real_snapshot_node = snapshot_filesystem_node

            def snapshot_node_with_fault(path: Path) -> tuple[str, ...]:
                if capture_fault["enabled"] and path == fixture:
                    raise PermissionError("forced fixture capture denial")
                return real_snapshot_node(path)

            class ForcedPhaseFailure(unittest.TestCase):
                def runTest(inner_self) -> None:
                    patcher = mock.patch.object(
                        sys.modules[__name__],
                        "snapshot_filesystem_node",
                        side_effect=snapshot_node_with_fault,
                    )
                    patcher.start()
                    inner_self.addCleanup(patcher.stop)
                    register_template_fill_safety_finalizer(
                        inner_self,
                        sources,
                        external_root=external_root,
                        repository_state=repository_state,
                    )
                    state_file.write_text('{"state":"after"}\n', encoding="utf-8")
                    cached_module.write_bytes(b"after bytecode!")
                    ignored_skill_log.write_text("ignored external write\n", encoding="utf-8")
                    (external_root / "unexpected.txt").write_text("dirty\n", encoding="utf-8")
                    capture_fault["enabled"] = True
                    inner_self.fail("forced primary phase failure")

            result = unittest.TestResult()
            ForcedPhaseFailure().run(result)
            diagnostics = "\n".join(
                detail
                for _, detail in [*result.failures, *result.errors]
            )
            self.assertIn("forced primary phase failure", diagnostics)
            self.assertIn("external ppt-master Git status", diagnostics)
            self.assertIn("external ppt-master __pycache__", diagnostics)
            self.assertIn("external ppt-master skill tree", diagnostics)
            self.assertIn("repository runtime .slidesmith", diagnostics)
            self.assertIn("source fixtures capture failed", diagnostics)


@unittest.skipUnless(
    os.environ.get("SLIDESMITH_RUN_REAL_TEMPLATE_FILL_SMOKE") == "1",
    "set SLIDESMITH_RUN_REAL_TEMPLATE_FILL_SMOKE=1 for the real Template Fill smoke",
)
class RealTemplateFillRunnerSmokeTests(unittest.TestCase):
    def test_check_apply_validate(self) -> None:
        fixture_value = os.environ.get("SLIDESMITH_TEMPLATE_FILL_SMOKE_PPTX", "").strip()
        content_value = os.environ.get("SLIDESMITH_TEMPLATE_FILL_SMOKE_CONTENT", "").strip()
        self.assertTrue(fixture_value, "SLIDESMITH_TEMPLATE_FILL_SMOKE_PPTX is required")
        self.assertTrue(content_value, "SLIDESMITH_TEMPLATE_FILL_SMOKE_CONTENT is required")

        fixture = Path(fixture_value).expanduser()
        content_source = Path(content_value).expanduser()
        expected_sources = {
            fixture: (105_402, "ddcaef381c298e0c5f4d1c636731044ed513e30166c07e79dae70ff5896227a3"),
            content_source: (12_239, "c823b8ff8e5733d7e0a778256cd0d03cf7b71da3cd8bb2cf15a70d5431d72906"),
        }
        external_root = ppt_runner.PPT_MASTER_ROOT
        repository_state = RUNNER_PATH.parent.parent / ".slidesmith"
        safety_snapshot = register_template_fill_safety_finalizer(
            self,
            list(expected_sources),
            external_root=external_root,
            repository_state=repository_state,
        )
        git_status = safety_snapshot["external ppt-master Git status"]
        self.assertEqual(git_status[0], 0, git_status)
        self.assertEqual(git_status[1], "", git_status)

        source_digests_before: dict[Path, str] = {}
        for path, (expected_size, expected_digest) in expected_sources.items():
            self.assertFalse(path.is_symlink(), f"smoke source must not be a symlink: {path}")
            self.assertTrue(path.is_file(), f"smoke source must be a regular file: {path}")
            self.assertEqual(path.stat().st_size, expected_size, path)
            source_digests_before[path] = sha256_file(path)
            self.assertEqual(source_digests_before[path], expected_digest, path)

        temporary = tempfile.TemporaryDirectory()
        self.addCleanup(temporary.cleanup)
        root = Path(temporary.name)
        workspace = root / "workspace"
        project = workspace / "projects" / "spec3_template_fill_smoke"
        try:
            sources_dir = project / "sources"
            analysis_dir = project / "analysis"
            sources_dir.mkdir(parents=True)
            analysis_dir.mkdir(parents=True)
            copied_fixture = sources_dir / fixture.name
            copied_content = sources_dir / content_source.name
            shutil.copy2(fixture, copied_fixture)
            shutil.copy2(content_source, copied_content)
            self.assertEqual(sha256_file(copied_fixture), source_digests_before[fixture])
            self.assertEqual(sha256_file(copied_content), source_digests_before[content_source])

            analyzer = ppt_runner.SCRIPTS_DIR / "template_fill_pptx.py"
            slide_library = analysis_dir / f"{copied_fixture.stem}.slide_library.json"
            analyze = run_template_fill_smoke_subprocess(
                [
                    sys.executable,
                    str(analyzer),
                    "analyze",
                    str(copied_fixture),
                    "-o",
                    str(slide_library),
                ],
                workspace=workspace,
                temporary_root=root,
            )
            self.assertEqual(analyze.returncode, 0, analyze.stdout + analyze.stderr)
            self.assertTrue(slide_library.is_file())
            self.assertGreater(slide_library.stat().st_size, 0)

            fill_plan = analysis_dir / "fill_plan.json"
            write_json_file(
                fill_plan,
                {
                    "schema": "template_fill_pptx_plan.v1",
                    "status": "confirmed",
                    "source_pptx": f"sources/{copied_fixture.name}",
                    "accepted_warnings": [],
                    "slides": [
                        {
                            "source_slide": 1,
                            "purpose": "automated Template Fill smoke cover",
                            "layout_rationale": {
                                "layout_pattern": "blueprint cover",
                                "why_fit": "A short source-grounded label fits the native cover slot.",
                                "risk": "Low; one short editable text replacement.",
                            },
                            "transition": "fade",
                            "replacements": [
                                {"slot_id": "s01_sh10", "text": "Cluster topology"}
                            ],
                            "table_edits": [],
                            "chart_edits": [],
                        }
                    ],
                },
            )
            plan = json.loads(fill_plan.read_text(encoding="utf-8"))
            self.assertEqual(plan["status"], "confirmed")
            self.assertEqual(len(plan["slides"]), 1)
            plan_digest = sha256_file(fill_plan)

            def run_phase(*arguments: str) -> subprocess.CompletedProcess[str]:
                completed = run_template_fill_smoke_subprocess(
                    [sys.executable, str(RUNNER_PATH), *arguments],
                    workspace=workspace,
                    temporary_root=root,
                )
                self.assertEqual(
                    completed.returncode,
                    0,
                    f"{' '.join(arguments)}\n{completed.stdout}{completed.stderr}",
                )
                return completed

            check_report = analysis_dir / "check_report.json"
            self.assertFalse(check_report.exists())
            run_phase("template-fill-check", "--project-path", str(project))
            self.assertEqual(sha256_file(fill_plan), plan_digest)
            self.assertTrue(check_report.is_file())
            check = json.loads(check_report.read_text(encoding="utf-8"))
            self.assertEqual(check.get("schema"), "template_fill_pptx_check.v1")
            self.assertEqual(check.get("summary", {}).get("error"), 0)
            check_digest = sha256_file(check_report)

            exports_dir = project / "exports"
            export_pattern = re.compile(
                r"^spec3_template_fill_smoke_template_fill_\d{8}_\d{6}\.pptx$"
            )
            exports_before = (
                {path.name for path in exports_dir.iterdir() if export_pattern.fullmatch(path.name)}
                if exports_dir.is_dir()
                else set()
            )
            run_phase(
                "template-fill-apply",
                "--project-path",
                str(project),
                "--transition",
                "fade",
            )
            self.assertEqual(sha256_file(fill_plan), plan_digest)
            self.assertEqual(sha256_file(check_report), check_digest)
            exports_after = {
                path.name
                for path in exports_dir.iterdir()
                if export_pattern.fullmatch(path.name)
            }
            new_exports = exports_after - exports_before
            self.assertEqual(len(new_exports), 1, sorted(new_exports))
            export = exports_dir / new_exports.pop()
            self.assertFalse(export.is_symlink())
            self.assertTrue(export.is_file())
            self.assertGreater(export.stat().st_size, 0)
            self.assertFalse(
                (exports_dir / "spec3_template_fill_smoke_template_fill.pptx").exists()
            )

            validate_report = project / "validation" / "validate_report.json"
            readback = project / "validation" / "readback.md"
            self.assertFalse(validate_report.exists())
            self.assertFalse(readback.exists())
            run_phase("template-fill-validate", "--project-path", str(project))
            self.assertEqual(sha256_file(fill_plan), plan_digest)
            self.assertEqual(sha256_file(check_report), check_digest)
            self.assertTrue(validate_report.is_file())
            validate = json.loads(validate_report.read_text(encoding="utf-8"))
            self.assertEqual(validate.get("schema"), "template_fill_pptx_validate.v1")
            self.assertEqual(validate.get("summary", {}).get("error"), 0)
            self.assertTrue(readback.is_file())
            self.assertGreater(readback.stat().st_size, 0)
            self.assertIn("Cluster topology", readback.read_text(encoding="utf-8"))

            with zipfile.ZipFile(export) as package:
                slide_entries = [
                    name
                    for name in package.namelist()
                    if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)
                ]
            final_plan = json.loads(fill_plan.read_text(encoding="utf-8"))
            self.assertEqual(len(slide_entries), len(final_plan["slides"]))
            self.assertEqual(len(slide_entries), 1)

            for forbidden in (project / "design_spec.md", project / "spec_lock.md"):
                self.assertFalse(forbidden.exists(), forbidden)
            for svg_dir in (project / "svg_output", project / "svg_final"):
                svg_files = (
                    [path for path in svg_dir.rglob("*") if path.suffix.casefold() == ".svg"]
                    if svg_dir.exists()
                    else []
                )
                self.assertFalse(svg_files, svg_files)

            state_dir = workspace / ".slidesmith"
            self.assertTrue((state_dir / "events.ndjson").is_file())
            self.assertTrue((state_dir / "status.json").is_file())
            assert_template_fill_runner_state_isolated(root, workspace)

            self.assertEqual(sha256_file(copied_fixture), source_digests_before[fixture])
            self.assertEqual(sha256_file(copied_content), source_digests_before[content_source])
            for path, digest in source_digests_before.items():
                self.assertFalse(path.is_symlink(), path)
                self.assertTrue(path.is_file(), path)
                self.assertEqual(sha256_file(path), digest, path)
            print(
                "real Template Fill smoke: "
                f"plan_sha256={plan_digest} "
                f"check_report_sha256={check_digest} "
                f"export={export.name} "
                f"export_bytes={export.stat().st_size} "
                f"validate_errors={validate['summary']['error']} "
                f"readback_bytes={readback.stat().st_size} "
                f"slides={len(slide_entries)}"
            )
        except BaseException:
            try:
                retained = Path(tempfile.mkdtemp(prefix="slidesmith-template-fill-smoke-failed-"))
                shutil.copytree(root, retained, dirs_exist_ok=True)
            except Exception as retention_error:
                print(
                    f"could not retain failed Template Fill smoke: {retention_error}",
                    file=sys.stderr,
                )
            else:
                print(f"retained failed Template Fill smoke at {retained}", file=sys.stderr)
            raise


@unittest.skipUnless(
    os.environ.get("SLIDESMITH_RUN_REAL_SOURCE_SMOKES") == "1",
    "set SLIDESMITH_RUN_REAL_SOURCE_SMOKES=1 for external converter smokes",
)
class RealSourceConversionSmokeTests(unittest.TestCase):
    def test_all_six_ooxml_extensions_convert_and_analyze(self) -> None:
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as tmp:
            workspace = Path(tmp) / "workspace"
            uploads = workspace / "uploads" / "all-six"
            uploads.mkdir(parents=True)
            base = Path(tmp) / "base.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "Six format conversion smoke"
            slide.placeholders[1].text = "Readable by the real PPT Master converter."
            presentation.save(base)

            manifest_files = []
            for extension, content_type in CONTENT_TYPES.items():
                source = uploads / f"deck_{extension[1:]}{extension}"
                rewrite_ooxml_main_content_type(base, source, content_type)
                manifest_files.append(
                    {
                        "name": source.name,
                        "upload_path": source.relative_to(workspace).as_posix(),
                    }
                )
            write_source_manifest(workspace, manifest_files)

            completed = run_runner(
                workspace,
                "prepare",
                "--project",
                "all_six",
                "--profile",
                "smoke",
                "--sources-manifest",
                ".slidesmith/source_inputs.json",
            )
            self.assertEqual(completed.returncode, 0, completed.stdout + completed.stderr)
            project = only_project(workspace)
            self.assertTrue((project / "analysis" / "source_profile.json").is_file())
            for extension in CONTENT_TYPES:
                stem = f"deck_{extension[1:]}"
                archived = project / "sources" / f"{stem}{extension}"
                self.assertTrue(archived.is_file(), archived)
                self.assertTrue((project / "sources" / f"{stem}.md").is_file())
                self.assertTrue((project / "sources" / f"{stem}.conversion_profile.json").is_file())
                self.assertTrue((project / "analysis" / f"{stem}.identity.json").is_file())
                self.assertTrue((project / "analysis" / f"{stem}.slide_library.json").is_file())
                expected_type = (
                    CONTENT_TYPES[".pptm"]
                    if extension in {".pptm", ".ppsm", ".potm"}
                    else CONTENT_TYPES[".pptx"]
                )
                self.assertEqual(read_main_content_type(archived), expected_type)
                if extension in {".pptx", ".pptm"}:
                    self.assertEqual(archived.read_bytes(), (uploads / archived.name).read_bytes())

    def test_manifest_text_priority_and_legacy_input_fallback(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            manifest_workspace = root / "manifest-workspace"
            uploads = manifest_workspace / "uploads" / "manifest"
            uploads.mkdir(parents=True)
            selected = uploads / "notes.text"
            selected.write_text("# Text Source Signal\nselected-manifest-text\n", encoding="utf-8")
            fallback = uploads / "fallback.md"
            fallback.write_text("# Wrong source\nmust-not-be-selected\n", encoding="utf-8")
            write_source_manifest(
                manifest_workspace,
                [{"name": selected.name, "upload_path": selected.relative_to(manifest_workspace).as_posix()}],
            )

            manifest_run = run_runner(
                manifest_workspace,
                "prepare",
                "--project",
                "manifest_text",
                "--profile",
                "real-lite",
                "--sources-manifest",
                ".slidesmith/source_inputs.json",
                "--input",
                fallback.relative_to(manifest_workspace).as_posix(),
            )
            self.assertEqual(manifest_run.returncode, 0, manifest_run.stdout + manifest_run.stderr)
            manifest_project = only_project(manifest_workspace)
            self.assertTrue((manifest_project / "sources" / "notes.text").is_file())
            self.assertFalse((manifest_project / "sources" / "fallback.md").exists())
            recommendations = json.loads(
                (manifest_project / "confirm_ui" / "recommendations.json").read_text(encoding="utf-8")
            )
            self.assertIn("Text Source Signal", recommendations["content_divergence"]["value"])

            legacy_workspace = root / "legacy-workspace"
            legacy = legacy_workspace / "uploads" / "legacy" / "legacy.md"
            legacy.parent.mkdir(parents=True)
            legacy.write_text("# Legacy Signal\nselected-legacy-input\n", encoding="utf-8")
            legacy_run = run_runner(
                legacy_workspace,
                "prepare",
                "--project",
                "legacy_input",
                "--profile",
                "smoke",
                "--input",
                legacy.relative_to(legacy_workspace).as_posix(),
            )
            self.assertEqual(legacy_run.returncode, 0, legacy_run.stdout + legacy_run.stderr)
            legacy_project = only_project(legacy_workspace)
            self.assertEqual(
                (legacy_project / "sources" / "legacy.md").read_text(encoding="utf-8"),
                legacy.read_text(encoding="utf-8"),
            )


def stage_one(source: Path, destination: Path) -> Path:
    destination.mkdir(parents=True)
    args = argparse.Namespace(sources_manifest="", input=str(source))
    staged = ppt_runner.stage_prepare_inputs(args, destination)
    if len(staged) != 1:
        raise AssertionError(f"staged inputs = {staged!r}, want exactly one")
    return staged[0]


def write_ooxml_package(
    path: Path,
    main_content_type: str,
    *,
    include_presentation: bool = True,
    presentation_xml: bytes | None = None,
    content_types_namespace: str = CONTENT_TYPES_NAMESPACE,
    duplicate_main_override: bool = False,
) -> None:
    override = f'  <Override PartName="/ppt/presentation.xml" ContentType="{main_content_type}"/>\n'
    overrides = override * (2 if duplicate_main_override else 1)
    content_types = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="{content_types_namespace}">
  <Default Extension="xml" ContentType="application/xml"/>
{overrides.rstrip()}
</Types>
'''.encode("utf-8")
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as package:
        package.writestr("[Content_Types].xml", content_types)
        if include_presentation:
            package.writestr(
                "ppt/presentation.xml",
                presentation_xml
                or f"<p:presentation xmlns:p='{PRESENTATIONML_NAMESPACE}'/>".encode("utf-8"),
            )
        package.writestr("ppt/marker.bin", b"preserve-me")


def rewrite_ooxml_main_content_type(source: Path, destination: Path, main_content_type: str) -> None:
    with zipfile.ZipFile(source) as source_package, zipfile.ZipFile(destination, "w") as destination_package:
        destination_package.comment = source_package.comment
        for entry in source_package.infolist():
            data = source_package.read(entry)
            if entry.filename == "[Content_Types].xml":
                original = CONTENT_TYPES[".pptx"].encode("utf-8")
                replacement = main_content_type.encode("utf-8")
                if data.count(original) != 1:
                    raise AssertionError("base PPTX does not contain exactly one presentation main content type")
                data = data.replace(original, replacement, 1)
            destination_package.writestr(entry, data)


def write_source_manifest(workspace: Path, files: list[dict[str, str]]) -> None:
    manifest = workspace / ".slidesmith" / "source_inputs.json"
    manifest.parent.mkdir(parents=True, exist_ok=True)
    manifest.write_text(
        json.dumps(
            {
                "schema": "slidesmith.source_inputs.v1",
                "task_id": "source-smoke",
                "files": files,
            },
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )


def run_runner(workspace: Path, *args: str) -> subprocess.CompletedProcess[str]:
    env = os.environ.copy()
    env["WORKSPACE"] = str(workspace)
    return subprocess.run(
        [sys.executable, str(RUNNER_PATH), *args],
        cwd=RUNNER_PATH.parent.parent,
        env=env,
        text=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        timeout=180,
    )


def only_project(workspace: Path) -> Path:
    projects = [path for path in (workspace / "projects").iterdir() if path.is_dir()]
    if len(projects) != 1:
        raise AssertionError(f"projects = {projects!r}, want exactly one")
    return projects[0]


def read_main_content_type(path: Path) -> str:
    raw = read_zip_entry(path, "[Content_Types].xml").decode("utf-8")
    marker = 'PartName="/ppt/presentation.xml" ContentType="'
    return raw.split(marker, 1)[1].split('"', 1)[0]


def read_zip_entry(path: Path, name: str) -> bytes:
    with zipfile.ZipFile(path) as package:
        return package.read(name)


if __name__ == "__main__":
    unittest.main()
